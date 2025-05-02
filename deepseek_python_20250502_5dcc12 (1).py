import tkinter as tk
from tkinter import scrolledtext, filedialog, messagebox, ttk, colorchooser
from tkinter import font
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import re

def sanitize_filename(name):
    name = re.sub(r'[\\/:*?"<>|]', '_', name)
    name = re.sub(r'\s+', '_', name)
    name = name.strip().strip('.')
    max_length = 200
    return name[:max_length] if len(name) > max_length else name

def is_valid_filename(filename):
    invalid_chars = set('\\/:*?"<>|')
    reserved_names = {'CON', 'PRN', 'AUX', 'NUL', 
                     'COM1', 'COM2', 'COM3', 'COM4', 'COM5', 'COM6', 'COM7', 'COM8', 'COM9',
                     'LPT1', 'LPT2', 'LPT3', 'LPT4', 'LPT5', 'LPT6', 'LPT7', 'LPT8', 'LPT9'}
    base_name = os.path.splitext(filename)[0].upper()
    return (
        not any(char in invalid_chars for char in filename) and
        len(filename) <= 255 and
        base_name not in reserved_names and
        not filename.startswith((' ', '.')) and
        not filename.endswith((' ', '.'))
    )

def create_gui():
    root = tk.Tk()
    root.title("TXT轉PPT")
    root.geometry("750x500")
    root.configure(bg="#E6F2FF")

    ming_font = font.Font(family="細明體", size=10)

    # 输入区域
    input_frame = tk.Frame(root, bg="#E6F2FF")
    input_frame.place(x=20, y=20)

    input_label = tk.Label(input_frame, text="輸入文字:", font=ming_font, bg="#E6F2FF")
    input_label.pack(anchor="w")

    input_text = scrolledtext.ScrolledText(input_frame, width=36, height=12, font=ming_font)
    input_text.pack(pady=5)

    # 右键菜单
    right_click_menu = tk.Menu(root, tearoff=0)
    right_click_menu.add_command(label="剪下", command=lambda: input_text.event_generate("<<Cut>>"))
    right_click_menu.add_command(label="複製", command=lambda: input_text.event_generate("<<Copy>>"))
    right_click_menu.add_command(label="貼上", command=lambda: input_text.event_generate("<<Paste>>"))
    right_click_menu.add_separator()
    right_click_menu.add_command(label="全選", command=lambda: input_text.event_generate("<<SelectAll>>"))

    def show_right_click_menu(event):
        right_click_menu.post(event.x_root, event.y_root)

    input_text.bind("<Button-3>", show_right_click_menu)

    # 预览区域
    preview_frame = tk.Frame(root, bg="#E6F2FF")
    preview_frame.place(x=400, y=20)

    preview_label = tk.Label(preview_frame, text="預覽:", font=ming_font, bg="#E6F2FF")
    preview_label.pack(anchor="w")

    preview_area = scrolledtext.ScrolledText(preview_frame, width=36, height=12, font=ming_font)
    preview_area.pack(pady=5)
    preview_area.config(state="disabled")

    # 按钮样式
    button_style = {"bg": "#4CAF50", "fg": "white", "font": ming_font, "padx": 10, "pady": 5}

    # 操作按钮框架
    button_frame = tk.Frame(root, bg="#E6F2FF")
    button_frame.place(x=20, y=300)

    # 设置区域
    settings_frame = tk.Frame(root, bg="#E6F2FF")
    settings_frame.place(x=20, y=350)

    # 全局变量
    global ppt_content, default_filename, selected_font, text_color, bg_color, slide_ratio
    ppt_content = ""
    default_filename = "簡報"
    selected_font = "細明體"
    text_color = (0, 0, 0)
    bg_color = (255, 255, 255)
    slide_ratio = "16:9"

    # ================== 核心功能函数定义（必须放在按钮创建之前） ==================
    def clear_text():
        input_text.delete(1.0, tk.END)
        update_preview("")

    def update_preview(text):
        preview_area.config(state="normal")
        preview_area.delete(1.0, tk.END)
        preview_area.insert(tk.END, text)
        preview_area.config(state="disabled")

    def process_text():
        text = input_text.get(1.0, tk.END).strip()
        if not text:
            messagebox.showwarning("警告", "請先輸入文字內容")
            return

        rule = page_rule_var.get()
        slides = []
        
        try:
            if rule == "自定义分页":
                separator = separator_var.get().encode().decode('unicode_escape')
                slides = [s.strip() for s in text.split(separator) if s.strip()]
            elif rule == "自动分页":
                threshold = int(auto_page_threshold.get())
                paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
                slides = [
                    "\n".join(paragraphs[i:i+threshold]) 
                    for i in range(0, len(paragraphs), threshold)
                ]
            elif rule == "多级分页":
                current_slide = []
                for line in text.splitlines():
                    line = line.strip()
                    if line.startswith("# "):
                        if current_slide:
                            slides.append("\n".join(current_slide))
                        current_slide = [line[2:]]
                    else:
                        current_slide.append(line)
                if current_slide:
                    slides.append("\n".join(current_slide))

            if not slides:
                messagebox.showwarning("警告", "未檢測到有效內容，請檢查輸入格式！")
                return

            preview_text = "將生成以下內容的PPT:\n\n"
            for i, slide_text in enumerate(slides[:5]):
                preview_text += f"第{i+1}頁: {slide_text.strip().splitlines()[0][:30]}...\n"

            if len(slides) > 5:
                preview_text += f"...共{len(slides)}頁"

            preview_text += f"\n\n字型: {selected_font}"
            preview_text += f"\n字體顏色: RGB{text_color}"
            preview_text += f"\n背景顏色: RGB{bg_color}"
            preview_text += f"\n版面比例: {slide_ratio}"

            update_preview(preview_text)

            raw_name = text[:4] if len(text) >= 4 else text
            filename = sanitize_filename(raw_name)
            if not filename.strip():
                filename = "簡報"
            global ppt_content, default_filename
            ppt_content = text
            default_filename = filename

            messagebox.showinfo("處理完成", f"文字已處理完成，請點擊儲存按鈕選擇保存位置\n預設檔名: {filename}")

        except ValueError:
            messagebox.showerror("錯誤", "自動分頁的段落數必須是整數")
            return

    def save_ppt():
        global ppt_content, default_filename, selected_font, text_color, bg_color, slide_ratio

        if not ppt_content.strip():
            messagebox.showerror("錯誤", "尚未處理任何文字內容，請先點擊「處理文字」")
            return

        file_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint 簡報", "*.pptx")],
            initialfile=default_filename
        )

        if not file_path:
            return

        try:
            filename = os.path.basename(file_path)
            if not is_valid_filename(filename):
                raise ValueError(f"文件名 '{filename}' 包含非法字符或為系統保留名稱")

            prs = Presentation()

            if slide_ratio == "16:9":
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
            elif slide_ratio == "4:3":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
            elif slide_ratio == "9:16":
                prs.slide_width = Inches(9)
                prs.slide_height = Inches(16)

            slide_w, slide_h = prs.slide_width, prs.slide_height
            box_w, box_h = int(slide_w * 0.85), int(slide_h * 0.85)
            left = int((slide_w - box_w) / 2)
            top = int((slide_h - box_h) / 2)

            slides = [s for s in ppt_content.split("\n\n") if s.strip()]
            for content in slides:
                lines = [line for line in content.splitlines() if line.strip()]
                if not lines:
                    continue

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bg_color)

                textbox = slide.shapes.add_textbox(left, top, box_w, box_h)
                text_frame = textbox.text_frame
                text_frame.clear()
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                title_p = text_frame.paragraphs[0]
                title_p.text = lines[0]
                title_p.alignment = PP_ALIGN.CENTER
                title_run = title_p.runs[0]
                title_run.font.name = selected_font
                title_run.font.size = Pt(48)
                title_run.font.color.rgb = RGBColor(*text_color)

                for line in lines[1:]:
                    p = text_frame.add_paragraph()
                    p.text = line
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]
                    run.font.name = selected_font
                    run.font.size = Pt(36)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(*text_color)

            prs.save(file_path)
            messagebox.showinfo("成功", f"PPT已成功儲存至:\n{file_path}")
        except IOError as e:
            messagebox.showerror("錯誤", f"路徑無效或權限不足:\n{str(e)}")
        except ValueError as e:
            messagebox.showerror("錯誤", f"文件名非法:\n{str(e)}")
        except Exception as e:
            messagebox.showerror("錯誤", f"儲存PPT時發生未知錯誤:\n{str(e)}")

    # ================== 界面控件定义 ==================
    # 分页规则设置
    page_rule_label = tk.Label(settings_frame, text="分页规则:", font=ming_font, bg="#E6F2FF")
    page_rule_label.grid(row=0, column=0, sticky="w", padx=5, pady=5)

    page_rule_var = tk.StringVar(value="自定义分页")
    page_rule_combo = ttk.Combobox(
        settings_frame, 
        textvariable=page_rule_var,
        values=["自定义分页", "自动分页", "多级分页"],
        width=12
    )
    page_rule_combo.grid(row=0, column=1, padx=5, pady=5)

    # 分页参数区域
    page_param_frame = tk.Frame(settings_frame, bg="#E6F2FF")
    page_param_frame.grid(row=0, column=2, columnspan=3, sticky="w")

    # 自定义分页组件
    separator_label = tk.Label(page_param_frame, text="分页符:", font=ming_font, bg="#E6F2FF")
    separator_var = tk.StringVar(value="\\n\\n")
    separator_entry = tk.Entry(page_param_frame, textvariable=separator_var, width=8)
    
    # 自动分页组件
    auto_page_label = tk.Label(page_param_frame, text="每页段落数:", font=ming_font, bg="#E6F2FF")
    auto_page_threshold = tk.Entry(page_param_frame, width=4)
    auto_page_threshold.insert(0, "3")
    
    # 多级分页组件
    multi_level_label = tk.Label(
        page_param_frame, 
        text="使用 # 分页，## 子标题",
        font=ming_font, 
        bg="#E6F2FF",
        fg="#666666"
    )

    def update_page_rule_ui(event=None):
        for widget in page_param_frame.winfo_children():
            widget.pack_forget()

        rule = page_rule_var.get()
        if rule == "自定义分页":
            separator_label.pack(side=tk.LEFT)
            separator_entry.pack(side=tk.LEFT, padx=5)
        elif rule == "自动分页":
            auto_page_label.pack(side=tk.LEFT)
            auto_page_threshold.pack(side=tk.LEFT, padx=5)
        elif rule == "多级分页":
            multi_level_label.pack(side=tk.LEFT)

    page_rule_combo.bind("<<ComboboxSelected>>", update_page_rule_ui)
    update_page_rule_ui()

    # 字体设置
    def is_chinese_font(font_name):
        chinese_keywords = ['細明體', '新細明體', '標楷體', '微軟', '黑體', '宋體', '華康', '文鼎', '文泉', '思源','麥克', 'DFKai', 'PMingLiU', 'MingLiU']
        return any(keyword in font_name for keyword in chinese_keywords)

    available_fonts = [f for f in font.families() if is_chinese_font(f)]
    font_var = tk.StringVar()
    font_var.set("細明體")

    font_combo = ttk.Combobox(settings_frame, textvariable=font_var, values=available_fonts, width=15)
    font_combo.grid(row=1, column=1, padx=5, pady=5)

    def on_font_change(event):
        global selected_font
        selected_font = font_var.get()

    font_combo.bind("<<ComboboxSelected>>", on_font_change)

    # 颜色选择
    def choose_text_color():
        global text_color
        color = colorchooser.askcolor(title="選擇字體顏色")
        if color[0]:
            text_color = (int(color[0][0]), int(color[0][1]), int(color[0][2]))
            text_color_button.config(bg=color[1])

    def choose_bg_color():
        global bg_color
        color = colorchooser.askcolor(title="選擇背景顏色")
        if color[0]:
            bg_color = (int(color[0][0]), int(color[0][1]), int(color[0][2]))
            bg_color_button.config(bg=color[1])

    text_color_button = tk.Button(settings_frame, text="字體顏色", command=choose_text_color, **button_style)
    text_color_button.grid(row=1, column=2, padx=5, pady=5)

    bg_color_button = tk.Button(settings_frame, text="背景顏色", command=choose_bg_color, **button_style)
    bg_color_button.grid(row=1, column=3, padx=5, pady=5)

    # 比例设置
    ratio_label = tk.Label(settings_frame, text="版面比例:", font=ming_font, bg="#E6F2FF")
    ratio_label.grid(row=2, column=0, sticky="w", padx=5, pady=5)

    ratio_var = tk.StringVar()
    ratio_var.set("16:9")

    ratio_combo = ttk.Combobox(settings_frame, textvariable=ratio_var, values=["16:9", "4:3", "9:16"], width=10)
    ratio_combo.grid(row=2, column=1, padx=5, pady=5)

    def on_ratio_change(event):
        global slide_ratio
        slide_ratio = ratio_var.get()

    ratio_combo.bind("<<ComboboxSelected>>", on_ratio_change)

    # 功能按钮
    process_button = tk.Button(button_frame, text="處理文字", command=process_text, **button_style)
    process_button.pack(side=tk.LEFT, padx=10)

    clear_button = tk.Button(button_frame, text="清除", command=clear_text, **button_style)
    clear_button.pack(side=tk.LEFT, padx=10)

    save_button = tk.Button(button_frame, text="儲存", command=save_ppt, **button_style)
    save_button.pack(side=tk.LEFT, padx=10)

    exit_button = tk.Button(button_frame, text="退出", command=root.destroy, **button_style)
    exit_button.pack(side=tk.LEFT, padx=10)

    root.mainloop()

if __name__ == "__main__":
    create_gui()