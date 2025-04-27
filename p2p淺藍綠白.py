from tkinterdnd2 import DND_FILES, TkinterDnD
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import threading
import time
import os
import tempfile
import sys
import subprocess
from typing import Optional

def check_poppler():
    try:
        from pdf2image.exceptions import PDFPageCountError
        convert_from_path(os.path.join(os.path.dirname(__file__), 'test.pdf'), first_page=1, last_page=1)
        return True
    except Exception as e:
        if 'poppler' in str(e).lower():
            messagebox.showerror(
                '錯誤',
                '請先安裝Poppler並確保其在系統PATH中。\n\n安裝步驟：\n'
                '1. 下載Poppler for Windows\n'
                '2. 解壓縮到C:\\Program Files\\poppler\n'
                '3. 將C:\\Program Files\\poppler\\bin加入系統PATH環境變量\n'
                '4. 重新啟動程式'
            )
            sys.exit(1)
        return True

check_poppler()

class PDFToPPTConverter:
    def __init__(self, root: TkinterDnD.Tk):
        self.root = root
        self.root.title("PDF 轉 PPT 轉換器")
        self.root.configure(bg="#ADD8E6")  # 淺藍色背景
        self.root.geometry("600x500")
        self.root.resizable(False, False)

        self.pdf_path: Optional[str] = None
        self.ppt_path: Optional[str] = None

        self.last_dir = os.path.expanduser("~/Desktop")

        self.converting = False
        self.current_page = 0
        self.total_pages = 0

        self._create_widgets()

        self.root.drop_target_register(DND_FILES)
        self.root.dnd_bind('<<Drop>>', self.on_drop)

    def _create_widgets(self):
        font_setting = ("微軟正黑體", 12)

        # --- PDF 輸入區 ---
        pdf_frame = tk.Frame(self.root, bg="#ADD8E6")
        pdf_frame.pack(pady=15)
        tk.Label(pdf_frame, text="PDF 文件:", bg="#ADD8E6", font=font_setting).pack(side=tk.LEFT)
        self.pdf_entry = tk.Entry(pdf_frame, width=40, font=font_setting)
        self.pdf_entry.pack(side=tk.LEFT, padx=10)
        tk.Button(pdf_frame, text="瀏覽", font=font_setting, command=self.select_pdf, relief="flat", bg="#4CAF50", fg="#FFFFFF").pack(side=tk.LEFT)

        # --- PPT 儲存區 ---
        ppt_frame = tk.Frame(self.root, bg="#ADD8E6")
        ppt_frame.pack(pady=15)
        tk.Label(ppt_frame, text="PPT 輸出:", bg="#ADD8E6", font=font_setting).pack(side=tk.LEFT)
        self.ppt_entry = tk.Entry(ppt_frame, width=40, font=font_setting)
        self.ppt_entry.pack(side=tk.LEFT, padx=10)
        tk.Button(ppt_frame, text="儲存", font=font_setting, command=self.select_save_location, relief="flat", bg="#4CAF50", fg="#FFFFFF").pack(side=tk.LEFT)

        # --- 比例選單 ---
        ratio_frame = tk.Frame(self.root, bg="#ADD8E6")
        ratio_frame.pack(pady=10)
        tk.Label(ratio_frame, text="選擇PPT比例:", bg="#ADD8E6", font=font_setting).pack()
        self.aspect_ratio = tk.StringVar(value="16:9")
        self.aspect_dropdown = ttk.Combobox(
            ratio_frame, textvariable=self.aspect_ratio,
            values=["16:9", "4:3", "10:16"], state="readonly", width=10, font=font_setting
        )
        self.aspect_dropdown.pack(pady=5)

        # --- 轉換按鈕 ---
        self.convert_btn = tk.Button(self.root, text="開始轉換", font=font_setting,
                                     command=self.start_conversion, relief="flat", width=20, height=2, bg="#4CAF50", fg="#FFFFFF")
        self.convert_btn.pack(pady=15)

        # --- 進度條 ---
        style = ttk.Style()
        style.configure("green.Horizontal.TProgressbar", troughcolor="#DDD", background="#4CAF50", thickness=20)
        self.progressbar = ttk.Progressbar(self.root, mode="determinate", style="green.Horizontal.TProgressbar")
        self.progressbar.pack(fill="x", padx=30, pady=10)

        # --- 動畫提示 ---
        self.anim_label = tk.Label(self.root, text="", bg="#ADD8E6", font=("微軟正黑體", 11))
        self.anim_label.pack(pady=5)

    def select_pdf(self):
        filename = filedialog.askopenfilename(
            filetypes=[("PDF 文件", "*.pdf")],
            initialdir=self.last_dir
        )
        if filename:
            self.pdf_path = filename
            self.last_dir = os.path.dirname(filename)
            self.pdf_entry.delete(0, tk.END)
            self.pdf_entry.insert(0, filename)
            self.save_ppt()

    def select_save_location(self):
        if not self.pdf_path:
            messagebox.showerror("錯誤", "請先選擇 PDF 文件！")
            return

        filename = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            initialfile=f"{os.path.splitext(os.path.basename(self.pdf_path))[0]}.pptx",  # 修正這裡
            filetypes=[("PowerPoint 文件", "*.pptx")],
            initialdir=self.last_dir
        )
        if filename:
            self.ppt_path = filename
            self.last_dir = os.path.dirname(filename)
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, filename)

    def save_ppt(self):
        if self.pdf_path:
            base = os.path.splitext(os.path.basename(self.pdf_path))[0]  # 修正這裡
            self.ppt_path = os.path.join(os.path.dirname(self.pdf_path), f"{base}.pptx")
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, self.ppt_path)

    def on_drop(self, event):
        path = event.data.strip('{}')
        if path.lower().endswith(".pdf"):
            self.pdf_path = path
            self.last_dir = os.path.dirname(path)
            self.pdf_entry.delete(0, tk.END)
            self.pdf_entry.insert(0, path)
            self.save_ppt()

    def start_conversion(self):
        if not self.pdf_path or not self.ppt_path:
            messagebox.showerror("錯誤", "請選擇輸入和輸出文件！")
            return

        self.convert_btn.config(state=tk.DISABLED)
        self.converting = True
        self.current_page = 0
        self.done_event = threading.Event()
        self.success_flag = {'ok': True}

        threading.Thread(target=self._run_conversion).start()
        threading.Thread(target=self._run_progress).start()

    def _run_conversion(self):
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                prs = Presentation()

                ratio = self.aspect_ratio.get()
                if ratio == "16:9":
                    prs.slide_width = Inches(16)
                    prs.slide_height = Inches(9)
                elif ratio == "4:3":
                    prs.slide_width = Inches(10)
                    prs.slide_height = Inches(7.5)
                elif ratio == "10:16":
                    prs.slide_width = Inches(9)
                    prs.slide_height = Inches(16)

                images = convert_from_path(self.pdf_path, fmt='jpeg', thread_count=4)
                self.total_pages = len(images)

                for i, image in enumerate(images):
                    self.current_page = i + 1
                    img_path = os.path.join(temp_dir, f'page_{i}.jpg')
                    image.save(img_path, 'JPEG')
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 修正這裡 (0 通常是標題投影片，6 是空白)
                    slide.shapes.add_picture(img_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)

                prs.save(self.ppt_path)

        except Exception as e:
            self.success_flag['ok'] = False
            self.root.after(0, lambda: self._show_error(str(e)))
        finally:
            self.done_event.set()

    def _run_progress(self):
        while not self.done_event.is_set():
            if self.total_pages > 0:
                percent = (self.current_page / self.total_pages) * 100
                self.progressbar["value"] = percent
            else:
                self.progressbar["value"] = 0
            self.root.update_idletasks()
            time.sleep(0.1)

        if self.success_flag['ok']:
            self.converting = False
            self.anim_label.config(text="轉換完成 ✔")
            self.root.after(1000, self._finish_success)

    def _finish_success(self):
        self.convert_btn.config(state=tk.NORMAL)
        messagebox.showinfo("成功", "文件轉換完成！")
        self.progressbar["value"] = 0
        self.anim_label.config(text="")

        # 自動開啟 PPT 所在資料夾
        folder = os.path.dirname(self.ppt_path)
        if sys.platform == "win32":
            os.startfile(folder)
        elif sys.platform == "darwin":
            subprocess.run(["open", folder])
        else:
            subprocess.run(["xdg-open", folder])

    def _show_error(self, msg: str):
        self.converting = False
        self.anim_label.config(text="轉換失敗 ✘")
        messagebox.showerror("錯誤", f"轉換失敗：{msg}")
        self.convert_btn.config(state=tk.NORMAL)
        self.progressbar["value"] = 0

if __name__ == "__main__":
    root = TkinterDnD.Tk()
    app = PDFToPPTConverter(root)
    root.mainloop()
