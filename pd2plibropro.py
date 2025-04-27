import tkinter as tk
from tkinter import ttk, filedialog, messagebox, colorchooser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import fitz  # PyMuPDF
import threading
import time
import os
import sys
import subprocess
from typing import Optional
from tkinter.font import families

class PDFToPPTConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF/Word 轉 PPT - 防呆版")
        self.root.configure(bg="#ADD8E6")
        self.root.geometry("450x360")
        self.root.resizable(False, False)

        self.input_path: Optional[str] = None
        self.ppt_path: Optional[str] = None
        self.input_type: Optional[str] = None
        self.last_dir = os.path.expanduser("~/Desktop")
        self.converting = False
        self.font_size = 24
        self.font_color = (0, 0, 0)
        self.font_name = "標楷體"
        self.page_bg_color = None

        self._create_widgets()

    def _create_widgets(self):
        font_setting = (self.font_name, 10)

        input_frame = tk.Frame(self.root, bg="#ADD8E6")
        input_frame.pack(pady=15)
        tk.Label(input_frame, text="輸入", bg="#ADD8E6", font=font_setting).pack(side=tk.LEFT)
        self.input_entry = tk.Entry(input_frame, width=30, font=font_setting)
        self.input_entry.pack(side=tk.LEFT, padx=10)
        tk.Button(input_frame, text="瀏覽", font=font_setting, command=self.select_file,
                  relief="flat", bg="#4CAF50", fg="#FFFFFF").pack(side=tk.LEFT)

        ppt_frame = tk.Frame(self.root, bg="#ADD8E6")
        ppt_frame.pack(pady=15)
        tk.Label(ppt_frame, text="結果", bg="#ADD8E6", font=font_setting).pack(side=tk.LEFT)
        self.ppt_entry = tk.Entry(ppt_frame, width=30, font=font_setting)
        self.ppt_entry.pack(side=tk.LEFT, padx=10)
        tk.Button(ppt_frame, text="另存", font=font_setting, command=self.select_save_location,
                  relief="flat", bg="#4CAF50", fg="#FFFFFF").pack(side=tk.LEFT)

        config_frame = tk.Frame(self.root, bg="#ADD8E6")
        config_frame.pack(pady=10)

        font_family_label = tk.Label(config_frame, text="字體", bg="#ADD8E6", font=font_setting)
        font_family_label.pack(side=tk.LEFT, padx=5)
        self.font_family_combobox = ttk.Combobox(config_frame, values=families(), font=font_setting, width=8)
        self.font_family_combobox.pack(side=tk.LEFT, padx=5)
        self.font_family_combobox.set("微軟正黑體")
        self.font_family_combobox.bind("<<ComboboxSelected>>", self.update_font_name)

        font_color_label = tk.Label(config_frame, text="字色", bg="#ADD8E6", font=font_setting)
        font_color_label.pack(side=tk.LEFT, padx=5)
        self.font_color_btn = tk.Button(config_frame, text="選色", command=self.choose_font_color,
                                        relief="flat", bg="#808080", fg="#FFFFFF", font=font_setting)
        self.font_color_btn.pack(side=tk.LEFT, padx=5)

        page_bg_color_label = tk.Label(config_frame, text="頁色", bg="#ADD8E6", font=font_setting)
        page_bg_color_label.pack(side=tk.LEFT, padx=5)
        self.page_bg_color_btn = tk.Button(config_frame, text="選色", command=self.choose_page_bg_color,
                                           relief="flat", bg="#808080", fg="#FFFFFF", font=font_setting)
        self.page_bg_color_btn.pack(side=tk.LEFT, padx=5)

        ratio_frame = tk.Frame(self.root, bg="#ADD8E6")
        ratio_frame.pack(pady=10)
        tk.Label(ratio_frame, text="比例", bg="#ADD8E6", font=font_setting).pack()
        self.aspect_ratio = tk.StringVar(value="16:9")
        self.aspect_dropdown = ttk.Combobox(ratio_frame, textvariable=self.aspect_ratio,
                                            values=["16:9", "4:3", "10:16"], state="readonly", width=6, font=font_setting)
        self.aspect_dropdown.pack(pady=5)

        self.convert_btn = tk.Button(self.root, text="開始轉檔", font=("標楷體", 18),
                                     command=self.start_conversion, relief="flat", width=8, height=1, bg="#4CAF50", fg="#FFFFDD")
        self.convert_btn.pack(pady=15)

        self.loading_label = tk.Label(self.root, text="", bg="#ADD8E6", font=("標楷體", 14))
        self.loading_label.pack(pady=5)

    def start_conversion(self):
        if not self.input_path or not self.ppt_path:
            messagebox.showerror("錯誤", "請選擇輸入和輸出文件！")
            return

        self.convert_btn.config(state=tk.DISABLED)
        self.converting = True
        self.success_flag = {'ok': True}

        if self.input_path.lower().endswith(".pdf"):
            self.input_type = "pdf"
        elif self.input_path.lower().endswith(".docx"):
            self.input_type = "docx"
        else:
            self.success_flag['ok'] = False
            self.root.after(0, lambda: self._show_error("不支援的檔案類型"))
            return

        threading.Thread(target=self._run_conversion).start()
        threading.Thread(target=self._animate_loading).start()

    def _animate_loading(self):
        dots = ""
        while self.converting:
            dots = dots + "." if len(dots) < 3 else ""
            self.loading_label.config(text=f"轉檔中{dots}")
            self.root.update_idletasks()
            time.sleep(0.5)

    def _run_conversion(self):
        try:
            prs = Presentation()
            ratio = self.aspect_ratio.get()
            if ratio == "16:9":
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
            elif ratio == "4:3":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
            elif ratio == "10:16":
                prs.slide_width = Inches(9)
                prs.slide_height = Inches(16)

            if self.input_type == 'docx':
                pdf_path = self._convert_docx_to_pdf()
                if not pdf_path:
                    return
            else:
                pdf_path = self.input_path

            doc = fitz.open(pdf_path)

            for page in doc:
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                if self.page_bg_color:
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*self.page_bg_color)

                textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
                tf = textbox.text_frame
                tf.clear()

                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = page.get_text("text").strip()

                run.font.size = Pt(self.font_size)
                run.font.color.rgb = RGBColor(*self.font_color)
                run.font.name = self.font_name

                p.alignment = PP_ALIGN.CENTER
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            prs.save(self.ppt_path)

            if self.input_type == 'docx' and os.path.exists(pdf_path):
                try:
                    os.remove(pdf_path)
                except PermissionError:
                    pass  # 如果檔案占用中，跳過刪除，不中斷程序

        except Exception as e:
            self.success_flag['ok'] = False
            self.root.after(0, lambda err=e: self._show_error(str(err)))
        finally:
            self.converting = False
            self.root.after(0, self._finish_success)

    def _convert_docx_to_pdf(self):
        pdf_path = os.path.splitext(self.input_path)[0] + ".pdf"
        try:
            subprocess.run([
                "soffice", "--headless", "--convert-to", "pdf", "--outdir", os.path.dirname(self.input_path), self.input_path
            ], check=True)
            return pdf_path
        except subprocess.CalledProcessError as e:
            self.success_flag['ok'] = False
            self.root.after(0, lambda err=e: self._show_error(f"DOCX 轉 PDF 失敗: {str(err)}"))
            return None

    def _finish_success(self):
        self.convert_btn.config(state=tk.NORMAL)
        if self.success_flag['ok']:
            messagebox.showinfo("成功", "文件轉換完成！")
        self.loading_label.config(text="")

    def select_file(self):
        filename = filedialog.askopenfilename(filetypes=[("所有支援格式", "*.pdf *.docx")], initialdir=self.last_dir)
        if filename:
            self.input_path = filename
            self.last_dir = os.path.dirname(filename)
            self.input_entry.delete(0, tk.END)
            self.input_entry.insert(0, filename)
            self.save_ppt()

    def select_save_location(self):
        if not self.input_path:
            messagebox.showerror("錯誤", "請先選擇輸入文件！")
            return
        filename = filedialog.asksaveasfilename(defaultextension=".pptx", initialfile=f"{os.path.splitext(os.path.basename(self.input_path))[0]}.pptx",
                                                 filetypes=[("PowerPoint 文件", "*.pptx")], initialdir=self.last_dir)
        if filename:
            self.ppt_path = filename
            self.last_dir = os.path.dirname(filename)
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, filename)

    def save_ppt(self):
        if self.input_path:
            base = os.path.splitext(os.path.basename(self.input_path))[0]
            self.ppt_path = os.path.join(os.path.dirname(self.input_path), f"{base}.pptx")
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, self.ppt_path)

    def update_font_name(self, event=None):
        self.font_name = self.font_family_combobox.get()

    def choose_font_color(self):
        color_code = colorchooser.askcolor(title="選擇字體顏色")[0]
        if color_code:
            self.font_color = (int(color_code[0]), int(color_code[1]), int(color_code[2]))
            self.font_color_btn.config(bg=f'#{int(color_code[0]):02x}{int(color_code[1]):02x}{int(color_code[2]):02x}')

    def choose_page_bg_color(self):
        color_code = colorchooser.askcolor(title="選擇頁面背景顏色")[0]
        if color_code:
            self.page_bg_color = (int(color_code[0]), int(color_code[1]), int(color_code[2]))
            self.page_bg_color_btn.config(bg=f'#{int(color_code[0]):02x}{int(color_code[1]):02x}{int(color_code[2]):02x}')

    def _show_error(self, msg: str):
        self.converting = False
        self.loading_label.config(text="轉換失敗 ✘")
        messagebox.showerror("錯誤", f"轉換失敗：{msg}")
        self.convert_btn.config(state=tk.NORMAL)

if __name__ == "__main__":
    root = tk.Tk()
    app = PDFToPPTConverter(root)
    root.mainloop()
