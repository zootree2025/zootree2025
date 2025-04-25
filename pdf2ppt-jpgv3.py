from tkinterdnd2 import DND_FILES, TkinterDnD
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import threading
import time
import os
import tempfile
import sys
from typing import Optional

def check_poppler():
    try:
        from pdf2image.exceptions import PDFPageCountError
        convert_from_path(os.path.join(os.path.dirname(__file__), 'test.pdf'), first_page=1, last_page=1)
        return True
    except Exception as e:
        if 'poppler' in str(e).lower():
            messagebox.showerror(
                '錯誤',
                '請先安裝Poppler並確保其在系統PATH中。\n\n安裝步驟：\n'
                '1. 下載Poppler for Windows\n'
                '2. 解壓縮到C:\\Program Files\\poppler\n'
                '3. 將C:\\Program Files\\poppler\\bin加入系統PATH環境變量\n'
                '4. 重新啟動程式'
            )
            sys.exit(1)
        return True

check_poppler()

class PDFToPPTConverter:
    def __init__(self, root: TkinterDnD.Tk):
        self.root = root
        self.root.title("PDF 轉 PPT 轉換器")
        self.root.configure(bg="Brown")

        self._create_widgets()

        self.pdf_path: Optional[str] = None
        self.ppt_path: Optional[str] = None

        self.converting = False
        self.current_page = 0
        self.total_pages = 0

    def _create_widgets(self):
        self._create_pdf_input_frame()
        self._create_ppt_output_frame()

        self.aspect_ratio = tk.StringVar(value="16:9")
        ttk.Label(self.root, text="選擇PPT比例:", background="brown", foreground="white", font=("Arial", 16)).pack()
        self.aspect_dropdown = ttk.Combobox(
            self.root, textvariable=self.aspect_ratio,
            values=["16:9", "4:3", "10:16"], font=("Arial", 16), state="readonly", width=10
        )
        self.aspect_dropdown.pack(pady=10)

        self.convert_btn = tk.Button(
            self.root, text="開始轉換", command=self.start_conversion,
            fg="white", bg="green", font=("Arial", 16), padx=20, pady=10
        )
        self.convert_btn.pack(pady=10)

        self.anim_label = tk.Label(
            self.root, text="", font=("Arial", 16), fg="white", bg="brown"
        )
        self.anim_label.pack()

        self.root.drop_target_register(DND_FILES)
        self.root.dnd_bind('<<Drop>>', self.on_drop)

    def _create_pdf_input_frame(self):
        self.pdf_frame = tk.Frame(self.root, bg="brown")
        self.pdf_frame.pack(padx=15, pady=15)
        tk.Label(self.pdf_frame, text="PDF 文件:", fg="white", bg="brown", font=("Arial", 16)).pack(side=tk.LEFT)
        self.pdf_entry = tk.Entry(self.pdf_frame, width=50, font=("Arial", 16))
        self.pdf_entry.pack(side=tk.LEFT, padx=15)
        tk.Button(
            self.pdf_frame, text="瀏覽", command=self.select_pdf,
            fg="white", bg="blue", font=("Arial", 16), padx=20, pady=10
        ).pack(side=tk.LEFT)

    def _create_ppt_output_frame(self):
        self.ppt_frame = tk.Frame(self.root, bg="brown")
        self.ppt_frame.pack(padx=15, pady=15)
        tk.Label(self.ppt_frame, text="PPT 輸出:", fg="white", bg="brown", font=("Arial", 16)).pack(side=tk.LEFT)
        self.ppt_entry = tk.Entry(self.ppt_frame, width=50, font=("Arial", 16))
        self.ppt_entry.pack(side=tk.LEFT, padx=15)
        tk.Button(
            self.ppt_frame, text="確認", command=self.save_ppt,
            fg="red", bg="Yellow", font=("Arial", 16), padx=20, pady=10
        ).pack(side=tk.LEFT)

    def on_drop(self, event):
        path = event.data.strip('{}')
        if path.lower().endswith(".pdf"):
            self.pdf_path = path
            self.pdf_entry.delete(0, tk.END)
            self.pdf_entry.insert(0, path)
            self.save_ppt()

    def select_pdf(self):
        filename = filedialog.askopenfilename(filetypes=[("PDF 文件", "*.pdf")])
        if filename:
            self.pdf_path = filename
            self.pdf_entry.delete(0, tk.END)
            self.pdf_entry.insert(0, filename)
            self.save_ppt()

    def save_ppt(self):
        if self.pdf_path:
            pdf_filename_without_ext = os.path.splitext(os.path.basename(self.pdf_path))[0]
            self.ppt_path = os.path.join(os.path.dirname(self.pdf_path), f"{pdf_filename_without_ext}.pptx")
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, self.ppt_path)
        else:
            messagebox.showerror("錯誤", "請先選擇 PDF 文件！")

    def start_conversion(self):
        if not self.pdf_path or not self.ppt_path:
            messagebox.showerror("錯誤", "請選擇輸入和輸出文件！")
            return

        self.convert_btn.config(state=tk.DISABLED)
        self.converting = True
        self.current_page = 0

        self.done_event = threading.Event()
        self.success_flag = {'ok': True}

        threading.Thread(target=self._run_conversion).start()
        threading.Thread(target=self._run_progress).start()

    def _run_conversion(self):
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                prs = Presentation()

                ratio = self.aspect_ratio.get()
                if ratio == "16:9":
                    prs.slide_width = Inches(16)
                    prs.slide_height = Inches(9)
                elif ratio == "4:3":
                    prs.slide_width = Inches(10)
                    prs.slide_height = Inches(7.5)
                elif ratio == "10:16":
                    prs.slide_width = Inches(9)
                    prs.slide_height = Inches(16)

                # 使用 multi-threaded 方式預先轉換圖片並快取
                images = convert_from_path(self.pdf_path, fmt='jpeg', thread_count=4)
                self.total_pages = len(images)

                for i, image in enumerate(images):
                    self.current_page = i + 1
                    self.root.after(10)
                    self.root.update_idletasks()
                    img_path = os.path.join(temp_dir, f'page_{i}.jpg')
                    image.save(img_path, 'JPEG')
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(img_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)

                prs.save(self.ppt_path)
        except Exception as e:
            error_msg = str(e)
            self.success_flag['ok'] = False
            self.root.after(0, lambda: self._show_error(error_msg))
        finally:
            self.done_event.set()

    def _run_progress(self):
        while not self.done_event.is_set():
            time.sleep(0.1)
            self.root.update_idletasks()
            if self.converting:
                self._update_anim_label()

        if self.success_flag['ok']:
            self.converting = False
            self.anim_label.config(text="轉換完成 ✔")
            self.root.after(1000, lambda: self._finish_success())

    def _update_anim_label(self):
        if self.converting:
            self.anim_label.config(text=f"...第 {self.current_page} 頁...")
            self.root.after(200, self._update_anim_label)

    def _finish_success(self):
        self.convert_btn.config(state=tk.NORMAL)
        messagebox.showinfo("成功", "文件轉換完成！")
        self.root.after(1500, lambda: self.anim_label.config(text=""))

    def _show_error(self, msg: str):
        self.converting = False
        self.anim_label.config(text="轉換失敗 ✘")
        messagebox.showerror("錯誤", f"轉換失敗：{msg}")
        self.convert_btn.config(state=tk.NORMAL)

if __name__ == "__main__":
    root = TkinterDnD.Tk()
    app = PDFToPPTConverter(root)
    root.mainloop()
