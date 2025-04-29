import tkinter as tk
from tkinter import ttk, filedialog, messagebox, colorchooser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import fitz  # PyMuPDF
import threading
import time
import os
import sys
import subprocess
from typing import Optional
from tkinter.font import families

class PDFToPPTConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF/Word 轉 PPT - 防呆版")
        self.root.configure(bg="#ADD8E6")
        self.root.geometry("450x360")
        self.root.resizable(False, False)

        self.input_path: Optional[str] = None
        self.ppt_path: Optional[str] = None
        self.input_type: Optional[str] = None
        self.last_dir = os.path.expanduser("~/Desktop")
        self.converting = False
        self.font_size = 24
        self.font_color = (0, 0, 0)
        self.font_name = "標楷體"
        self.page_bg_color = None
        self.text_align = tk.StringVar(value="LEFT")
        # 在 __init__ 方法添加
        self.template_path: Optional[str] = None

        self._create_widgets()

    def _create_widgets(self):
        font_setting = (self.font_name, 10)

        input_frame = tk.Frame(self.root, bg="#ADD8E6")
        input_frame.pack(pady=15)
        tk.Label(input_frame, text="輸入", bg="#ADD8E6", font=font_setting).pack(side=tk.LEFT)
        self.input_entry = tk.Entry(input_frame, width=30, font=font_setting)
        self.input_entry.pack(side=tk.LEFT, padx=10)
        tk.Button(input_frame, text="瀏覽", font=font_setting, command=self.select_file,
                  relief="flat", bg="#4CAF50", fg="#FFFFFF").pack(side=tk.LEFT)
        
        # 添加直接貼上文字按鈕
        tk.Button(input_frame, text="貼上文字", font=font_setting, command=self.open_text_input,
                  relief="flat", bg="#4CAF50", fg="#FFFFFF").pack(side=tk.LEFT, padx=5)

        ppt_frame = tk.Frame(self.root, bg="#ADD8E6")
        ppt_frame.pack(pady=15)
        tk.Label(ppt_frame, text="結果", bg="#ADD8E6", font=font_setting).pack(side=tk.LEFT)
        self.ppt_entry = tk.Entry(ppt_frame, width=30, font=font_setting)
        self.ppt_entry.pack(side=tk.LEFT, padx=10)
        tk.Button(ppt_frame, text="另存", font=font_setting, command=self.select_save_location,
                  relief="flat", bg="#4CAF50", fg="#FFFFFF").pack(side=tk.LEFT)

        config_frame = tk.Frame(self.root, bg="#ADD8E6")
        config_frame.pack(pady=10)

        font_family_label = tk.Label(config_frame, text="字體", bg="#ADD8E6", font=font_setting)
        font_family_label.pack(side=tk.LEFT, padx=5)
        self.font_family_combobox = ttk.Combobox(config_frame, values=families(), font=font_setting, width=8)
        self.font_family_combobox.pack(side=tk.LEFT, padx=5)
        self.font_family_combobox.set("微軟正黑體")
        self.font_family_combobox.bind("<<ComboboxSelected>>", self.update_font_name)

        font_color_label = tk.Label(config_frame, text="字色", bg="#ADD8E6", font=font_setting)
        font_color_label.pack(side=tk.LEFT, padx=5)
        self.font_color_btn = tk.Button(config_frame, text="選色", command=self.choose_font_color,
                                        relief="flat", bg="#808080", fg="#FFFFFF", font=font_setting)
        self.font_color_btn.pack(side=tk.LEFT, padx=5)

        page_bg_color_label = tk.Label(config_frame, text="頁色", bg="#ADD8E6", font=font_setting)
        page_bg_color_label.pack(side=tk.LEFT, padx=5)
        self.page_bg_color_btn = tk.Button(config_frame, text="選色", command=self.choose_page_bg_color,
                                           relief="flat", bg="#808080", fg="#FFFFFF", font=font_setting)
        self.page_bg_color_btn.pack(side=tk.LEFT, padx=5)

        # 添加文字對齊選項
        align_label = tk.Label(config_frame, text="對齊", bg="#ADD8E6", font=font_setting)
        align_label.pack(side=tk.LEFT, padx=5)
        self.align_dropdown = ttk.Combobox(config_frame, textvariable=self.text_align,
                                          values=["LEFT", "CENTER", "RIGHT"], state="readonly", width=6, font=font_setting)
        self.align_dropdown.pack(side=tk.LEFT, padx=5)

       
        ratio_frame = tk.Frame(self.root, bg="#ADD8E6")
        ratio_frame.pack(pady=10)
        tk.Label(ratio_frame, text="比例", bg="#ADD8E6", font=font_setting).pack(side=tk.LEFT)
        self.aspect_ratio = tk.StringVar(value="16:9")
        self.aspect_dropdown = ttk.Combobox(ratio_frame, textvariable=self.aspect_ratio,
                                            values=["16:9", "4:3", "10:16"], state="readonly", width=6, font=font_setting)
        self.aspect_dropdown.pack(side=tk.LEFT, padx=5)

        # 添加字體大小下拉選單
        font_size_frame = tk.Frame(ratio_frame, bg="#ADD8E6")
        font_size_frame.pack(side=tk.LEFT, padx=10)
        tk.Label(font_size_frame, text="字體大小", bg="#ADD8E6", font=font_setting).pack(side=tk.LEFT)
        self.font_size_var = tk.StringVar(value="24")
        font_sizes = [str(i) for i in range(8, 74, 2)]  # 8到72的偶數
        self.font_size_dropdown = ttk.Combobox(font_size_frame, textvariable=self.font_size_var,
                                               values=font_sizes, state="readonly", width=4, font=font_setting)
        self.font_size_dropdown.pack(side=tk.LEFT, padx=5)
        self.font_size_dropdown.bind("<<ComboboxSelected>>", self.update_font_size)

        # 添加模板選擇到字體大小右側
        template_label = tk.Label(font_size_frame, text="模板", bg="#ADD8E6", font=font_setting)
        template_label.pack(side=tk.LEFT, padx=5)
        self.template_combobox = ttk.Combobox(font_size_frame, values=["空白模板", "自定義..."], 
                                        state="readonly", width=8, font=font_setting)
        self.template_combobox.pack(side=tk.LEFT, padx=5)
        self.template_combobox.set("空白模板")
        self.template_combobox.bind("<<ComboboxSelected>>", self.update_template)
        
        # 模板瀏覽按鈕
        self.template_browse_btn = tk.Button(font_size_frame, text="瀏覽", command=self.select_template,
                                        relief="flat", bg="#808080", fg="#FFFFFF", font=font_setting)
        self.template_browse_btn.pack(side=tk.LEFT, padx=5)

        self.convert_btn = tk.Button(self.root, text="開始轉檔", font=("標楷體", 18),
                                     command=self.start_conversion, relief="flat", width=8, height=1, bg="#4CAF50", fg="#FFFFDD")
        self.convert_btn.pack(pady=15)

        self.loading_label = tk.Label(self.root, text="", bg="#ADD8E6", font=("標楷體", 14))
        self.loading_label.pack(pady=5)

    def start_conversion(self):
        if not self.input_path or not self.ppt_path:
            messagebox.showerror("錯誤", "請選擇輸入和輸出文件！")
            return

        self.convert_btn.config(state=tk.DISABLED)
        self.converting = True
        self.success_flag = {'ok': True}

        if self.input_path.lower().endswith(".pdf"):
            self.input_type = "pdf"
        elif self.input_path.lower().endswith(".docx"):
            self.input_type = "docx"
        elif self.input_path.lower().endswith(".txt"):
            self.input_type = "txt"
        else:
            self.success_flag['ok'] = False
            self.root.after(0, lambda: self._show_error("不支援的檔案類型"))
            return

        threading.Thread(target=self._run_conversion).start()
        threading.Thread(target=self._animate_loading).start()

    def _animate_loading(self):
        dots = ""
        while self.converting:
            dots = dots + "." if len(dots) < 3 else ""
            self.loading_label.config(text=f"轉檔中{dots}")
            self.root.update_idletasks()
            time.sleep(0.5)

    def _run_conversion(self):
        pdf_path = None
        doc = None
        try:
            # 根據模板創建簡報
            if self.template_path and os.path.exists(self.template_path):
                prs = Presentation(self.template_path)
            else:
                prs = Presentation()
            
            # 移除模板原有內容（根據需要保留部分內容）
            while len(prs.slides) > 0:
                xml_slides = prs.slides._sldIdLst  
                xml_slides.remove(xml_slides[0])
            
            ratio = self.aspect_ratio.get()
            if ratio == "16:9":
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
            elif ratio == "4:3":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
            elif ratio == "10:16":
                prs.slide_width = Inches(9)
                prs.slide_height = Inches(16)

            # 處理文字文件
            if self.input_path.endswith(".txt"):
                self.input_type = "txt"
                with open(self.input_path, "r", encoding="utf-8") as f:
                    text = f.read()
                
                # 根據兩個以上空白行分割文字內容
                page_contents = []
                current_page = []
                empty_line_count = 0
                
                for line in text.splitlines():
                    if not line.strip():  # 空行
                        empty_line_count += 1
                    else:
                        if empty_line_count >= 2:  # 兩個以上空白行表示換頁
                            if current_page:  # 確保不添加空頁面
                                page_contents.append("\n".join(current_page))
                                current_page = []
                        empty_line_count = 0
                        current_page.append(line)
                
                # 添加最後一頁
                if current_page:
                    page_contents.append("\n".join(current_page))
                
                # 如果沒有內容，創建一個空頁面
                if not page_contents:
                    page_contents = [""]
                
                # 為每一頁創建幻燈片
                for page_content in page_contents:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    if self.page_bg_color:
                        background = slide.background
                        fill = background.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(*self.page_bg_color)
                    
                    textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
                    tf = textbox.text_frame
                    tf.clear()
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    textbox.left = Inches(0)
                    textbox.top = Inches(0)
                    textbox.width = prs.slide_width
                    textbox.height = prs.slide_height
                    
                    lines = page_content.splitlines()
                    
                    for line in lines:
                        indent = 0
                        for char in line:
                            if char == ' ':
                                indent += 1
                            else:
                                break
                        level = min(max(indent // 4, 0), 8)
                        clean_line = line.lstrip()
                        
                        if not clean_line:
                            continue
                        
                        p = tf.add_paragraph()
                        p.text = clean_line
                        p.level = level
                        
                        align_value = self.text_align.get()
                        if align_value == "LEFT":
                            p.alignment = PP_ALIGN.LEFT
                        elif align_value == "CENTER":
                            p.alignment = PP_ALIGN.CENTER
                        elif align_value == "RIGHT":
                            p.alignment = PP_ALIGN.RIGHT
                        
                        run = p.runs[0]
                        run.font.size = Pt(self.font_size)
                        run.font.color.rgb = RGBColor(*self.font_color)
                        run.font.name = self.font_name
                    
                    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                
            elif self.input_type == 'docx':
                pdf_path = self._convert_docx_to_pdf()
                if not pdf_path:
                    return
            else:
                pdf_path = self.input_path
            
            # 處理 PDF 文件
            if self.input_type in ['pdf', 'docx']:
                doc = fitz.open(pdf_path)
                for page in doc:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])

                if self.page_bg_color:
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*self.page_bg_color)

                textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
                tf = textbox.text_frame
                tf.clear()
                # Set text frame to center align both horizontally and vertically
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center
                textbox.left = Inches(0)
                textbox.top = Inches(0)
                textbox.width = prs.slide_width
                textbox.height = prs.slide_height

                text = page.get_text("text", flags=0)
                lines = text.splitlines()

                for line in lines:
                    indent = 0
                    for char in line:
                        if char == ' ':
                            indent += 1
                        else:
                            break
                    level = min(max(indent // 4, 0), 8)  # 限制層級在 0~8 之間 
                    clean_line = line.lstrip()
                    
                    if not clean_line:  # 跳過空行
                        continue
                    
                    p = tf.add_paragraph()
                    p.text = clean_line
                    p.level = level
                    # 根據選擇的對齊方式設定文字對齊
                    align_value = self.text_align.get()
                    if align_value == "LEFT":
                        p.alignment = PP_ALIGN.LEFT
                    elif align_value == "CENTER":
                        p.alignment = PP_ALIGN.CENTER
                    elif align_value == "RIGHT":
                        p.alignment = PP_ALIGN.RIGHT

                    run = p.runs[0]
                    run.font.size = Pt(self.font_size)
                    run.font.color.rgb = RGBColor(*self.font_color)
                    run.font.name = self.font_name

                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            prs.save(self.ppt_path)

        except Exception as e:
            self.success_flag['ok'] = False
            self.root.after(0, lambda err=e: self._show_error(str(err)))
        finally:
            self.converting = False
            if doc:
                doc.close()
            if self.input_type == 'docx' and pdf_path and os.path.exists(pdf_path):
                try:
                    os.remove(pdf_path)
                except Exception as e:
                    print(f"刪除臨時檔案失敗: {e}")
            self.root.after(0, self._finish_success)

    def _convert_docx_to_pdf(self):
        try:
            pdf_path = os.path.splitext(self.input_path)[0] + ".pdf"
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", os.path.dirname(self.input_path), self.input_path],
                capture_output=True,
                text=True
            )
            if result.returncode != 0:
                raise Exception(f"LibreOffice 錯誤：\n{result.stderr}")
            return pdf_path
        except FileNotFoundError:
            self.root.after(0, lambda: self._show_error("請先安裝 LibreOffice"))
            return None
        except Exception as e:
            self.root.after(0, lambda: self._show_error(f"DOCX 轉 PDF 失敗: {str(e)}"))
            return None

    def _finish_success(self):
        self.convert_btn.config(state=tk.NORMAL)
        if self.success_flag['ok']:
            messagebox.showinfo("成功", "文件轉換完成！")
        self.loading_label.config(text="")

    def select_file(self):
        filename = filedialog.askopenfilename(filetypes=[("所有支援格式", "*.pdf *.docx")], initialdir=self.last_dir)
        if filename:
            self.input_path = filename
            self.last_dir = os.path.dirname(filename)
            self.input_entry.delete(0, tk.END)
            self.input_entry.insert(0, filename)
            self.save_ppt()

    def select_save_location(self):
        if not self.input_path:
            messagebox.showerror("錯誤", "請先選擇輸入文件！")
            return
        filename = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            initialfile=f"{os.path.splitext(os.path.basename(self.input_path))[0]}.pptx",
            filetypes=[("PowerPoint 文件", "*.pptx")],
            initialdir=self.last_dir
        )
        if filename:
            self.ppt_path = filename
            self.last_dir = os.path.dirname(filename)
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, filename)

    def save_ppt(self):
        if self.input_path:
            base = os.path.splitext(os.path.basename(self.input_path))[0]
            template_suffix = "_template" if self.template_path else ""
            self.ppt_path = os.path.join(os.path.dirname(self.input_path), 
                                       f"{base}{template_suffix}.pptx")
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, self.ppt_path)

    def update_font_name(self, event=None):
        self.font_name = self.font_family_combobox.get()

    def update_font_size(self, event=None):
        self.font_size = int(self.font_size_var.get())

    def choose_font_color(self):
        color_code = colorchooser.askcolor(title="選擇字體顏色")
        if color_code[0] is not None:
            rgb = color_code[0]
            if isinstance(rgb, tuple) and len(rgb) == 3:
                self.font_color = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
                self.font_color_btn.config(
                    bg=f'#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}')

    def choose_page_bg_color(self):
        color_code = colorchooser.askcolor(title="選擇頁面背景顏色")
        if color_code[0] is not None:
            rgb = color_code[0]
            if isinstance(rgb, tuple) and len(rgb) == 3:
                self.page_bg_color = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
                self.page_bg_color_btn.config(
                    bg=f'#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}')

    def _show_error(self, msg: str):
        self.converting = False
        self.loading_label.config(text="轉換失敗 ✘")
        messagebox.showerror("錯誤", f"轉換失敗：{msg}")
        self.convert_btn.config(state=tk.NORMAL)

    # 新增模板選擇方法
    def select_template(self):
        filename = filedialog.askopenfilename(
            filetypes=[("PPT模板", "*.pptx")],
            initialdir=self.last_dir
        )
        if filename:
            self.template_path = filename
            self.template_combobox.set(os.path.basename(filename))
    
    def update_template(self, event=None):
        if self.template_combobox.get() == "空白模板":
            self.template_path = None
            self.template_browse_btn.config(state=tk.NORMAL)
    
    # 添加文字輸入視窗方法
    def open_text_input(self):
        # 創建新視窗
        text_window = tk.Toplevel(self.root)
        text_window.title("輸入文字")
        text_window.geometry("800x600")  # 設置更大的初始尺寸
        text_window.configure(bg="#ADD8E6")
        
        # 設置視窗最大化
        text_window.state('zoomed')  # 在 Windows 上最大化視窗
        
        # 添加說明標籤
        tk.Label(text_window, text="請輸入或貼上文字內容：", bg="#ADD8E6", font=(self.font_name, 12)).pack(pady=10)
        
        # 添加文字框
        text_frame = tk.Frame(text_window, bg="#ADD8E6")
        text_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        text_box = tk.Text(text_frame, font=(self.font_name, 12), wrap=tk.WORD)
        text_box.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        # 添加滾動條
        scrollbar = tk.Scrollbar(text_frame, command=text_box.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        text_box.config(yscrollcommand=scrollbar.set)
        
        # 添加按鈕框架
        button_frame = tk.Frame(text_window, bg="#ADD8E6")
        button_frame.pack(pady=15)
        
        # 添加確認按鈕
        tk.Button(button_frame, text="確認", font=(self.font_name, 12), command=lambda: self.process_text_input(text_box.get("1.0", tk.END), text_window),
                  relief="flat", bg="#4CAF50", fg="#FFFFFF", width=8).pack(side=tk.LEFT, padx=10)
        
        # 添加取消按鈕
        tk.Button(button_frame, text="取消", font=(self.font_name, 12), command=text_window.destroy,
                  relief="flat", bg="#808080", fg="#FFFFFF", width=8).pack(side=tk.LEFT, padx=10)
    
    def process_text_input(self, text, window):
        if not text.strip():
            messagebox.showerror("錯誤", "請輸入文字內容！")
            return
        
        # 創建臨時文字文件
        temp_dir = os.path.join(os.path.expanduser("~"), "AppData", "Local", "Temp")
        os.makedirs(temp_dir, exist_ok=True)
        temp_file = os.path.join(temp_dir, "temp_text_input.txt")
        
        with open(temp_file, "w", encoding="utf-8") as f:
            f.write(text)
        
        # 設置輸入路徑為臨時文件
        self.input_path = temp_file
        self.input_type = "txt"  # 明確設置檔案類型
        self.input_entry.delete(0, tk.END)
        self.input_entry.insert(0, "直接輸入的文字")
        
        # 確保輸出目錄存在
        output_dir = os.path.dirname(os.path.abspath(__file__))  # 使用當前程式所在目錄
        
        # 設置輸出路徑
        self.ppt_path = os.path.join(output_dir, "文字輸入.pptx")
        self.ppt_entry.delete(0, tk.END)
        self.ppt_entry.insert(0, self.ppt_path)
        
        # 關閉視窗
        window.destroy()
        
        # 提示用戶
        messagebox.showinfo("成功", "文字已準備好，請點擊「開始轉檔」按鈕進行轉換。")

if __name__ == "__main__":
    root = tk.Tk()
    app = PDFToPPTConverter(root)
    root.mainloop()
