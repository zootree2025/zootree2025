from flask import Flask, request, jsonify, send_from_directory
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os
import tempfile
import uuid

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

@app.route('/convert', methods=['POST'])
def convert_pdf_to_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'Invalid file type, only PDFs are allowed'}), 400

    # Save uploaded file
    filename = f"{uuid.uuid4().hex}.pdf"
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(file_path)

    try:
        # Convert PDF to PPT
        with tempfile.TemporaryDirectory() as temp_dir:
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)

            images = convert_from_path(file_path, fmt='jpeg', thread_count=4)
            for i, image in enumerate(images):
                img_path = os.path.join(temp_dir, f'page_{i}.jpg')
                image.save(img_path, 'JPEG')
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)

            output_filename = f"{uuid.uuid4().hex}.pptx"
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            prs.save(output_path)

        return jsonify({'download_url': f'/download/{output_filename}'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        os.remove(file_path)

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
