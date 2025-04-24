import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import threading
import time
import os
import tempfile
import sys

def check_poppler():
    try:
        from pdf2image.exceptions import PDFPageCountError
        # 嘗試轉換一個空白PDF來測試poppler
        convert_from_path(os.path.join(os.path.dirname(__file__), 'test.pdf'))
        return True
    except Exception as e:
        if 'poppler' in str(e).lower():
            messagebox.showerror('錯誤', '請先安裝Poppler並確保其在系統PATH中。\n\n安裝步驟：\n1. 下載Poppler for Windows\n2. 解壓縮到C:\\Program Files\\poppler\n3. 將C:\\Program Files\\poppler\\bin加入系統PATH環境變量\n4. 重新啟動程式')
            sys.exit(1)
        return True

# 在程式啟動時檢查poppler
check_poppler()

class PDFToPPTConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF 轉 PPT 轉換器")

        # PDF 輸入
        self.pdf_frame = tk.Frame(self.root)
        self.pdf_frame.pack(padx=10, pady=5)
        tk.Label(self.pdf_frame, text="PDF 文件:").pack(side=tk.LEFT)
        self.pdf_entry = tk.Entry(self.pdf_frame, width=50)
        self.pdf_entry.pack(side=tk.LEFT, padx=5)
        tk.Button(self.pdf_frame, text="瀏覽", command=self.select_pdf).pack(side=tk.LEFT)

        # PPT 輸出
        self.ppt_frame = tk.Frame(self.root)
        self.ppt_frame.pack(padx=10, pady=5)
        tk.Label(self.ppt_frame, text="PPT 輸出:").pack(side=tk.LEFT)
        self.ppt_entry = tk.Entry(self.ppt_frame, width=50)
        self.ppt_entry.pack(side=tk.LEFT, padx=5)
        tk.Button(self.ppt_frame, text="瀏覽", command=self.save_ppt).pack(side=tk.LEFT)

        # 進度條
        self.progress = ttk.Progressbar(self.root, orient=tk.HORIZONTAL, length=300, mode='determinate')
        self.progress.pack(pady=10)

        # 轉換按鈕
        self.convert_btn = tk.Button(self.root, text="開始轉換", command=self.start_conversion)
        self.convert_btn.pack(pady=10)

        # 動畫標籤
        self.anim_label = tk.Label(self.root, text="", font=("Arial", 10), fg="blue")
        self.anim_label.pack()

    def select_pdf(self):
        filename = filedialog.askopenfilename(filetypes=[("PDF 文件", "*.pdf")])
        if filename:
            self.pdf_entry.delete(0, tk.END)
            self.pdf_entry.insert(0, filename)

    def save_ppt(self):
        filename = filedialog.asksaveasfilename(
            filetypes=[("PowerPoint 簡報", "*.pptx")],
            defaultextension=".pptx"
        )
        if filename:
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, filename)

    def start_conversion(self):
        pdf_path = self.pdf_entry.get()
        ppt_path = self.ppt_entry.get()

        if not pdf_path or not ppt_path:
            messagebox.showerror("錯誤", "請選擇輸入和輸出文件！")
            return

        self.convert_btn.config(state=tk.DISABLED)
        self.progress['value'] = 0

        filename = os.path.basename(pdf_path)
        anim_text = tk.StringVar()
        self.anim_label.config(textvariable=anim_text)

        done_event = threading.Event()
        success_flag = {'ok': True}

        def run_conversion():
            try:
                # 創建臨時目錄存放圖片
                with tempfile.TemporaryDirectory() as temp_dir:
                    # 轉換PDF為圖片
                    anim_text.set(f"正在轉換PDF為圖片：{filename}")
                    images = convert_from_path(pdf_path)
                    total_pages = len(images)

                    # 創建新的PPT
                    prs = Presentation()
                    
                    # 設置幻燈片大小為16:9
                    prs.slide_width = Inches(16)
                    prs.slide_height = Inches(9)

                    # 添加每一頁到PPT
                    for i, image in enumerate(images):
                        # 更新進度
                        progress = int((i / total_pages) * 90)
                        self.progress['value'] = progress
                        anim_text.set(f"正在處理第 {i+1}/{total_pages} 頁")

                        # 保存圖片到臨時文件
                        img_path = os.path.join(temp_dir, f'page_{i}.jpg')
                        image.save(img_path, 'JPEG')

                        # 添加幻燈片
                        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面配置
                        
                        # 設置圖片大小以適應幻燈片
                        left = top = Inches(0)
                        width = prs.slide_width
                        height = prs.slide_height
                        
                        # 添加圖片到幻燈片
                        slide.shapes.add_picture(img_path, left, top, width, height)

                    # 保存PPT
                    anim_text.set("正在保存PPT文件...")
                    prs.save(ppt_path)

            except Exception as e:
                error_msg = str(e)
                success_flag['ok'] = False
                self.root.after(0, lambda: self.show_error(error_msg))
            finally:
                done_event.set()

        def run_progress():
            while not done_event.is_set():
                time.sleep(0.05)
                self.root.update_idletasks()

            if success_flag['ok']:
                self.progress['value'] = 100
                anim_text.set("轉換完成 ✔")
                self.root.after(0, lambda: self.finish_success())

        threading.Thread(target=run_conversion).start()
        threading.Thread(target=run_progress).start()

    def finish_success(self):
        self.convert_btn.config(state=tk.NORMAL)
        messagebox.showinfo("成功", "文件轉換完成！")
        self.root.after(1500, lambda: self.anim_label.config(text=""))
        self.progress['value'] = 0

    def show_error(self, msg):
        self.anim_label.config(text="轉換失敗 ✘")
        messagebox.showerror("錯誤", f"轉換失敗：{msg}")
        self.convert_btn.config(state=tk.NORMAL)
        self.progress['value'] = 0

if __name__ == "__main__":
    root = tk.Tk()
    app = PDFToPPTConverter(root)
    root.mainloop()
