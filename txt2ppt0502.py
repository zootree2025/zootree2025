import tkinter as tk
from tkinter import scrolledtext, filedialog, messagebox, ttk, colorchooser
from tkinter import font
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_gui():
    # 創建主視窗
    root = tk.Tk()
    root.title("TXT轉PPT")
    root.geometry("650x450")  # 增加高度以容納新控件
    root.configure(bg="#E6F2FF")  # 淡藍色背景
    
    # 設定字體
    ming_font = font.Font(family="細明體", size=10)
    
    # 創建左側輸入文字窗
    input_frame = tk.Frame(root, bg="#E6F2FF")
    input_frame.place(x=20, y=20)
    
    input_label = tk.Label(input_frame, text="輸入文字:", font=ming_font, bg="#E6F2FF")
    input_label.pack(anchor="w")
    
    input_text = scrolledtext.ScrolledText(input_frame, width=36, height=12, font=ming_font)
    input_text.pack(pady=5)
    
    # 創建右鍵選單
    right_click_menu = tk.Menu(root, tearoff=0)
    right_click_menu.add_command(label="剪下", command=lambda: input_text.event_generate("<<Cut>>"))
    right_click_menu.add_command(label="複製", command=lambda: input_text.event_generate("<<Copy>>"))
    right_click_menu.add_command(label="貼上", command=lambda: input_text.event_generate("<<Paste>>"))
    right_click_menu.add_separator()
    right_click_menu.add_command(label="全選", command=lambda: input_text.event_generate("<<SelectAll>>"))
    
    # 綁定右鍵選單
    def show_right_click_menu(event):
        right_click_menu.post(event.x_root, event.y_root)
    
    input_text.bind("<Button-3>", show_right_click_menu)
    
    # 創建右側預覽視窗
    preview_frame = tk.Frame(root, bg="#E6F2FF")
    preview_frame.place(x=330, y=20)
    
    preview_label = tk.Label(preview_frame, text="預覽:", font=ming_font, bg="#E6F2FF")
    preview_label.pack(anchor="w")
    
    preview_area = scrolledtext.ScrolledText(preview_frame, width=36, height=12, font=ming_font)
    preview_area.pack(pady=5)
    preview_area.config(state="disabled")  # 預覽區域設為只讀
    
    # 創建按鈕樣式 (綠底白字)
    button_style = {"bg": "#4CAF50", "fg": "white", "font": ming_font, "padx": 10, "pady": 5}
    
    # 底部按鈕區域
    button_frame = tk.Frame(root, bg="#E6F2FF")
    button_frame.place(x=20, y=280)
    
    # 新增設定區域
    settings_frame = tk.Frame(root, bg="#E6F2FF")
    settings_frame.place(x=20, y=330)
    
    settings_label = tk.Label(settings_frame, text="PPT設定:", font=ming_font, bg="#E6F2FF")
    settings_label.grid(row=0, column=0, sticky="w", padx=5, pady=5)
    
    # 新增全域變數
    global ppt_content, default_filename, selected_font, text_color, bg_color, slide_ratio
    ppt_content = ""
    default_filename = "簡報.pptx"
    selected_font = "細明體"
    text_color = (0, 0, 0)  # 黑色
    bg_color = (255, 255, 255)  # 白色
    slide_ratio = "16:9"  # 預設比例
    
    # 字型選擇功能
    def on_font_change(event):
        global selected_font
        selected_font = font_var.get()
    
    # 獲取系統可用字型並篩選中文字型
    def is_chinese_font(font_name):
        # 常見中文字型關鍵字
        chinese_keywords = ['細明體', '新細明體', '標楷體', '微軟', '黑體', '宋體', 
                          '華康', '文鼎', '文泉', '思源','麥克', 'DFKai', 'PMingLiU', 'MingLiU']
        return any(keyword in font_name for keyword in chinese_keywords)
    
    available_fonts = [f for f in font.families() if is_chinese_font(f)]
    font_var = tk.StringVar()
    font_var.set("細明體")  # 預設值
    
    # 字體顏色選擇功能
    def choose_text_color():
        global text_color
        color = colorchooser.askcolor(title="選擇字體顏色")
        if color[0]:  # 如果用戶選擇了顏色而不是取消
            text_color = (int(color[0][0]), int(color[0][1]), int(color[0][2]))
            text_color_button.config(bg=color[1])
    
    # 背景顏色選擇功能
    def choose_bg_color():
        global bg_color
        color = colorchooser.askcolor(title="選擇背景顏色")
        if color[0]:  # 如果用戶選擇了顏色而不是取消
            bg_color = (int(color[0][0]), int(color[0][1]), int(color[0][2]))
            bg_color_button.config(bg=color[1])
    
    # 清除按鈕功能
    def clear_text():
        input_text.delete(1.0, tk.END)
        update_preview("")
    
    # 更新預覽區域
    def update_preview(text):
        preview_area.config(state="normal")
        preview_area.delete(1.0, tk.END)
        preview_area.insert(tk.END, text)
        preview_area.config(state="disabled")
    
    # 處理文字功能
    def process_text():
        text = input_text.get(1.0, tk.END).strip()
        if not text:
            messagebox.showwarning("警告", "請先輸入文字內容")
            return
        
        # 更新預覽
        preview_text = "將生成以下內容的PPT:\n\n"
        lines = text.split('\n')
        for i, line in enumerate(lines[:5]):
            preview_text += f"第{i+1}頁: {line[:30]}{'...' if len(line) > 30 else ''}\n"
        
        if len(lines) > 5:
            preview_text += f"...共{len(lines)}頁"
        
        # 添加設定信息到預覽
        preview_text += f"\n\n字型: {selected_font}"
        preview_text += f"\n字體顏色: RGB{text_color}"
        preview_text += f"\n背景顏色: RGB{bg_color}"
        preview_text += f"\n版面比例: {slide_ratio}"
        
        update_preview(preview_text)
        
        # 生成檔名 (使用前四個字)
        if len(text) >= 4:
            filename = text[:4] + ".pptx"
        else:
            filename = text + ".pptx"
        
        # 儲存全域變數以便儲存功能使用
        global ppt_content, default_filename
        ppt_content = text
        default_filename = filename
        
        messagebox.showinfo("處理完成", f"文字已處理完成，請點擊儲存按鈕選擇保存位置\n預設檔名: {filename}")
    
    # 儲存功能
    def save_ppt():
        global ppt_content, default_filename, selected_font, text_color, bg_color, slide_ratio
        
        if not hasattr(save_ppt, 'ppt_content') and not 'ppt_content' in globals():
            messagebox.showwarning("警告", "請先處理文字")
            return
        
        # 選擇儲存位置
        file_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint 簡報", "*.pptx")],
            initialfile=default_filename
        )
        
        if not file_path:
            return  # 用戶取消儲存
        
        try:
            # 創建PPT
            prs = Presentation()
            
            # 設定投影片比例
            if slide_ratio == "16:9":
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
            elif slide_ratio == "4:3":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
            elif slide_ratio == "9:16":
                prs.slide_width = Inches(9)
                prs.slide_height = Inches(16)
            
            # 分行處理文字，每行一張投影片
            lines = ppt_content.split('\n')
            for line in lines:
                if line.strip():  # 跳過空行
                    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 使用標題和內容版面配置
                    
                    # 設定背景顏色
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])
                    
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    # 設定標題文字
                    title.text = line[:50]  # 標題最多50字
                    
                    # 設定標題字型、顏色和置中
                    for paragraph in title.text_frame.paragraphs:
                        paragraph.alignment = 1  # 1 代表置中對齊
                        for run in paragraph.runs:
                            run.font.name = selected_font
                            run.font.color.rgb = RGBColor(text_color[0], text_color[1], text_color[2])
                    
                    # 設定內容文字
                    if len(line) > 50:
                        content.text = line[50:]  # 剩餘內容
                        
                        # 設定內容字型、顏色和置中
                        for paragraph in content.text_frame.paragraphs:
                            paragraph.alignment = 1  # 1 代表置中對齊
                            for run in paragraph.runs:
                                run.font.name = selected_font
                                run.font.color.rgb = RGBColor(text_color[0], text_color[1], text_color[2])
            
            # 儲存PPT
            prs.save(file_path)
            messagebox.showinfo("成功", f"PPT已成功儲存至:\n{file_path}")
        except Exception as e:
            messagebox.showerror("錯誤", f"儲存PPT時發生錯誤:\n{str(e)}")
    
    # 按鈕
    process_button = tk.Button(button_frame, text="處理文字", command=process_text, **button_style)
    process_button.pack(side=tk.LEFT, padx=10)
    
    clear_button = tk.Button(button_frame, text="清除", command=clear_text, **button_style)
    clear_button.pack(side=tk.LEFT, padx=10)
    
    save_button = tk.Button(button_frame, text="儲存", command=save_ppt, **button_style)
    save_button.pack(side=tk.LEFT, padx=10)
    
    exit_button = tk.Button(button_frame, text="退出", command=root.destroy, **button_style)
    exit_button.pack(side=tk.LEFT, padx=10)
    
    # 新增設定按鈕
    font_combo = ttk.Combobox(settings_frame, textvariable=font_var, values=available_fonts, width=15)
    font_combo.grid(row=0, column=1, padx=5, pady=5)
    font_combo.bind("<<ComboboxSelected>>", on_font_change)
    
    text_color_button = tk.Button(settings_frame, text="字體顏色", command=choose_text_color, **button_style)
    text_color_button.grid(row=0, column=2, padx=5, pady=5)
    
    bg_color_button = tk.Button(settings_frame, text="背景顏色", command=choose_bg_color, **button_style)
    bg_color_button.grid(row=0, column=3, padx=5, pady=5)
    
    # 比例下拉選單
    ratio_label = tk.Label(settings_frame, text="版面比例:", font=ming_font, bg="#E6F2FF")
    ratio_label.grid(row=1, column=0, sticky="w", padx=5, pady=5)
    
    ratio_var = tk.StringVar()
    ratio_var.set("16:9")  # 預設值
    
    def on_ratio_change(event):
        global slide_ratio
        slide_ratio = ratio_var.get()
        messagebox.showinfo("比例設定", f"已設定版面比例為: {slide_ratio}")
    
    ratio_combo = ttk.Combobox(settings_frame, textvariable=ratio_var, values=["16:9", "4:3", "9:16"], width=10)
    ratio_combo.grid(row=1, column=1, padx=5, pady=5)
    ratio_combo.bind("<<ComboboxSelected>>", on_ratio_change)
    
    # 啟動主循環
    root.mainloop()

if __name__ == "__main__":
    create_gui()